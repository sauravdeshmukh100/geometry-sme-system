"""
Complete Geometry Document Processor with Integrated OCR Support

OCR Trigger Logic:
==================
1. PDF Files: 
   - First tries native PyPDF text extraction
   - If NO text found (empty/scanned PDF) → Triggers OCR
   
2. Image Files (.png, .jpg, .jpeg, .tiff, .bmp):
   - ALWAYS triggers OCR (no alternative)
   
3. Other Formats (DOCX, PPTX, TXT, MD, HTML):
   - Never use OCR (have native text extraction)
"""

import os
import json
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime
import logging
import re
import io

# Document processing libraries
import pypdf
from docx import Document as DocxDocument
from pptx import Presentation
import markdown
from bs4 import BeautifulSoup

# OCR libraries
import fitz  # PyMuPDF
from PIL import Image
import pytesseract
import cv2
import numpy as np

from ..config.settings import settings

logger = logging.getLogger(__name__)


class GeometryOCRProcessor:
    """
    Advanced OCR processor for geometry textbooks.
    Handles text extraction, diagram detection, and formula recognition.
    
    This class is called when:
    1. PDF has no native text (scanned/image-based PDF)
    2. Processing standalone image files (.png, .jpg, etc.)
    """
    
    def __init__(
        self,
        tesseract_path: Optional[str] = None,
        min_confidence: float = 60.0,
        enable_preprocessing: bool = True
    ):
        """
        Initialize OCR processor.
        
        Args:
            tesseract_path: Path to tesseract executable (if not in PATH)
            min_confidence: Minimum OCR confidence threshold (0-100)
            enable_preprocessing: Enable image preprocessing for better OCR
        """
        if tesseract_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_path
        
        self.min_confidence = min_confidence
        self.enable_preprocessing = enable_preprocessing
        
        # Geometry-specific patterns for content identification
        self.diagram_keywords = [
            'figure', 'diagram', 'illustration', 'image',
            'triangle', 'circle', 'rectangle', 'polygon',
            'angle', 'line', 'point', 'vertex'
        ]
        
        self.formula_patterns = [
            r'[A-Z]\s*=\s*',  # Variable assignments (A = ...)
            r'\b(sin|cos|tan|log|sqrt)\s*\(',  # Trig functions
            r'\d+\s*[+\-*/]\s*\d+',  # Arithmetic operations
            r'[a-z]\^[0-9]',  # Exponents (a^2)
            r'\\frac\{',  # LaTeX fractions
            r'π|∏|Σ|∫|√|°',  # Math symbols
        ]
        
        logger.info(f"OCR Processor initialized (min_confidence={min_confidence}%, "
                   f"preprocessing={enable_preprocessing})")
    
    def process_pdf_with_ocr(
        self,
        pdf_path: str,
        extract_images: bool = True,
        ocr_images: bool = True
    ) -> Dict[str, Any]:
        """
        Process a PDF with OCR capabilities.
        
        ⚠️ CALLED WHEN: PDF has no native text (scanned/image-based PDF)
        
        Args:
            pdf_path: Path to PDF file
            extract_images: Extract images from PDF
            ocr_images: Apply OCR to images
        
        Returns:
            Dictionary with extracted text and metadata
        """
        logger.info(f"🔍 OCR Processing started for: {pdf_path}")
        
        results = {
            'file_path': pdf_path,
            'total_pages': 0,
            'pages': [],
            'full_text': '',
            'has_diagrams': False,
            'has_formulas': False,
            'image_count': 0,
            'avg_confidence': 0.0,
            'ocr_used': True  # This function only called when OCR is needed
        }
        
        try:
            doc = fitz.open(pdf_path)
            results['total_pages'] = len(doc)
            
            confidences = []
            
            for page_num in range(len(doc)):
                logger.info(f"  📄 OCR processing page {page_num + 1}/{len(doc)}")
                page = doc[page_num]
                
                # Extract native text first (might have some text even in scanned PDFs)
                native_text = page.get_text()
                
                # Extract images and apply OCR
                page_images = []
                image_texts = []
                
                if extract_images:
                    images = self._extract_page_images(page, page_num)
                    results['image_count'] += len(images)
                    logger.info(f"    🖼️  Found {len(images)} images on page {page_num + 1}")
                    
                    if ocr_images and images:
                        for img_index, img_data in enumerate(images):
                            logger.info(f"      🔎 Running OCR on image {img_index + 1}/{len(images)}")
                            ocr_result = self._ocr_image(
                                img_data['image'],
                                img_data['bbox']
                            )
                            if ocr_result:
                                image_texts.append(ocr_result['text'])
                                confidences.append(ocr_result['confidence'])
                                page_images.append({
                                    'bbox': img_data['bbox'],
                                    'text': ocr_result['text'],
                                    'confidence': ocr_result['confidence']
                                })
                                logger.info(f"        ✓ OCR confidence: {ocr_result['confidence']:.1f}%")
                            else:
                                logger.warning(f"        ✗ OCR failed for image {img_index + 1}")
                
                # Combine native text and OCR text
                combined_text = native_text
                if image_texts:
                    combined_text += "\n\n[IMAGE CONTENT]\n" + "\n".join(image_texts)
                
                # Analyze page content
                contains_diagrams = self._detect_diagrams(combined_text)
                contains_formulas = self._detect_formulas(combined_text)
                
                results['pages'].append({
                    'page_num': page_num + 1,
                    'text': combined_text,
                    'has_images': len(page_images) > 0,
                    'contains_diagrams': contains_diagrams,
                    'contains_formulas': contains_formulas,
                    'confidence': np.mean(confidences) if confidences else 100.0
                })
                
                results['full_text'] += combined_text + "\n\n"
                
                if contains_diagrams:
                    results['has_diagrams'] = True
                if contains_formulas:
                    results['has_formulas'] = True
            
            doc.close()
            
            if confidences:
                results['avg_confidence'] = np.mean(confidences)
                logger.info(f"✅ OCR complete: {results['total_pages']} pages, "
                          f"{results['image_count']} images, "
                          f"avg confidence: {results['avg_confidence']:.2f}%")
            else:
                logger.warning(f"⚠️  OCR completed but no text extracted from {pdf_path}")
            
        except Exception as e:
            logger.error(f"❌ Error in OCR processing {pdf_path}: {e}", exc_info=True)
            raise
        
        return results
    
    def _extract_page_images(
        self,
        page: fitz.Page,
        page_num: int
    ) -> List[Dict[str, Any]]:
        """Extract images from a PDF page."""
        images = []
        
        try:
            image_list = page.get_images(full=True)
            
            for img_index, img_info in enumerate(image_list):
                try:
                    xref = img_info[0]
                    base_image = page.parent.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Convert to PIL Image
                    image = Image.open(io.BytesIO(image_bytes))
                    
                    # Get image location on page
                    img_rects = page.get_image_rects(xref)
                    bbox = img_rects[0] if img_rects else None
                    
                    images.append({
                        'image': image,
                        'bbox': bbox,
                        'index': img_index,
                        'page': page_num + 1
                    })
                    
                except Exception as e:
                    logger.warning(f"Failed to extract image {img_index} from page {page_num}: {e}")
        
        except Exception as e:
            logger.warning(f"Failed to get images from page {page_num}: {e}")
        
        return images
    
    def _preprocess_image(self, image: Image.Image) -> np.ndarray:
        """
        Preprocess image for better OCR results.
        
        Steps:
        1. Convert to grayscale
        2. Apply denoising
        3. Adaptive thresholding
        4. Dilation to connect broken characters
        """
        # Convert PIL to OpenCV format
        img_array = np.array(image)
        
        # Convert to grayscale
        if len(img_array.shape) == 3:
            gray = cv2.cvtColor(img_array, cv2.COLOR_RGB2GRAY)
        else:
            gray = img_array
        
        # Apply denoising
        denoised = cv2.fastNlMeansDenoising(gray)
        
        # Apply adaptive thresholding
        thresh = cv2.adaptiveThreshold(
            denoised,
            255,
            cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
            cv2.THRESH_BINARY,
            11,
            2
        )
        
        # Dilation to connect broken characters
        kernel = np.ones((1, 1), np.uint8)
        dilated = cv2.dilate(thresh, kernel, iterations=1)
        
        return dilated
    
    def _ocr_image(
        self,
        image: Image.Image,
        bbox: Optional[fitz.Rect] = None
    ) -> Optional[Dict[str, Any]]:
        """
        Apply OCR to an image.
        
        Returns:
            Dictionary with text and confidence, or None if OCR fails
        """
        try:
            # Preprocess if enabled
            if self.enable_preprocessing:
                processed = self._preprocess_image(image)
                image = Image.fromarray(processed)
            
            # Run OCR with detailed output
            # PSM 6 = Assume a single uniform block of text
            ocr_data = pytesseract.image_to_data(
                image,
                output_type=pytesseract.Output.DICT,
                config='--psm 6'
            )
            
            # Extract text and calculate confidence
            texts = []
            confidences = []
            
            for i, text in enumerate(ocr_data['text']):
                conf = int(ocr_data['conf'][i])
                if conf > self.min_confidence and text.strip():
                    texts.append(text)
                    confidences.append(conf)
            
            if not texts:
                return None
            
            combined_text = ' '.join(texts)
            avg_confidence = np.mean(confidences) if confidences else 0.0
            
            return {
                'text': combined_text,
                'confidence': avg_confidence,
                'bbox': bbox
            }
            
        except Exception as e:
            logger.warning(f"OCR failed for image: {e}")
            return None
    
    def _detect_diagrams(self, text: str) -> bool:
        """Detect if text mentions diagrams/figures."""
        text_lower = text.lower()
        return any(keyword in text_lower for keyword in self.diagram_keywords)
    
    def _detect_formulas(self, text: str) -> bool:
        """Detect mathematical formulas in text."""
        for pattern in self.formula_patterns:
            if re.search(pattern, text):
                return True
        return False
    
    def process_image_file(self, image_path: str) -> Dict[str, Any]:
        """
        Process a standalone image file.
        
        ⚠️ CALLED WHEN: Processing .png, .jpg, .jpeg, .tiff, .bmp files
        
        Args:
            image_path: Path to image file
        
        Returns:
            Dictionary with OCR results
        """
        logger.info(f"🔍 OCR Processing image: {image_path}")
        
        try:
            image = Image.open(image_path)
            
            ocr_result = self._ocr_image(image, None)
            
            if ocr_result:
                logger.info(f"  ✅ OCR confidence: {ocr_result['confidence']:.2f}%")
                return {
                    'file_path': image_path,
                    'text': ocr_result['text'],
                    'confidence': ocr_result['confidence'],
                    'contains_diagrams': self._detect_diagrams(ocr_result['text']),
                    'contains_formulas': self._detect_formulas(ocr_result['text']),
                    'ocr_used': True
                }
            else:
                logger.warning(f"  ⚠️  OCR returned no text for {image_path}")
                return {
                    'file_path': image_path,
                    'text': '',
                    'confidence': 0.0,
                    'error': 'OCR failed - no text detected',
                    'ocr_used': True
                }
        
        except Exception as e:
            logger.error(f"❌ Error processing image {image_path}: {e}")
            return {
                'file_path': image_path,
                'text': '',
                'confidence': 0.0,
                'error': str(e),
                'ocr_used': False
            }


class GeometryDocumentProcessor:
    """
    Enhanced document processor with OCR support for geometry content.
    
    OCR Decision Flow:
    ==================
    PDF Files:
        1. Try PyPDF native text extraction
        2. Check if text is empty/whitespace only
        3. If empty → Call OCR processor
        4. If has text → Use native extraction
    
    Image Files (.png, .jpg, etc.):
        → Always call OCR processor
    
    Other Files (DOCX, PPTX, TXT, MD, HTML):
        → Never use OCR (have native text extraction)
    """
    
    SUPPORTED_FORMATS = {'.pdf', '.docx', '.pptx', '.txt', '.md', '.html', 
                         '.png', '.jpg', '.jpeg', '.tiff', '.bmp'}
    
    # Geometry-specific patterns
    GEOMETRY_PATTERNS = {
        'theorem': r'\b(theorem|proof|lemma|corollary|proposition)\b',
        'shape': r'\b(triangle|square|rectangle|circle|polygon|quadrilateral|pentagon|hexagon|octagon)\b',
        'angle': r'\b(angle|degree|radian|acute|obtuse|right angle|complementary|supplementary)\b',
        'formula': r'\b(area|perimeter|volume|circumference|pythagorean|surface area)\b',
        'concept': r'\b(congruent|similar|parallel|perpendicular|bisect|symmetry|transformation)\b'
    }
    
    # Grade-specific keywords for classification
    GRADE_LEVEL_KEYWORDS = {
        'grade_6': ['basic shapes', 'perimeter', 'area of rectangle', 'area of square', 
                    'circle basics', 'simple angles', 'line segment', 'ray'],
        'grade_7': ['triangles', 'congruence', 'construction', 'symmetry', 
                    'area of parallelogram', 'area of triangle'],
        'grade_8': ['quadrilaterals', 'mensuration', 'volume', 'surface area',
                    'understanding quadrilaterals', 'visualizing solid shapes'],
        'grade_9': ['coordinate geometry', 'lines and angles', 'triangles advanced',
                    'heron formula', 'areas of parallelograms', 'circles'],
        'grade_10': ['similar triangles', 'pythagoras theorem', 'trigonometry',
                     'circles advanced', 'tangent', 'constructions', 'areas related to circles']
    }
    
    def __init__(self, enable_ocr: bool = True, ocr_config: Optional[Dict[str, Any]] = None):
        """
        Initialize document processor with OCR support.
        
        Args:
            enable_ocr: Enable OCR for scanned PDFs and images
            ocr_config: OCR configuration dict with keys:
                - tesseract_path: Path to tesseract (optional)
                - min_confidence: Minimum confidence threshold (default: 60.0)
                - enable_preprocessing: Enable image preprocessing (default: True)
        """
        self.processed_docs = []
        self.enable_ocr = enable_ocr
        
        # Initialize OCR processor if enabled
        if self.enable_ocr:
            ocr_config = ocr_config or {}
            self.ocr_processor = GeometryOCRProcessor(
                tesseract_path=ocr_config.get('tesseract_path'),
                min_confidence=ocr_config.get('min_confidence', 60.0),
                enable_preprocessing=ocr_config.get('enable_preprocessing', True)
            )
            logger.info("✅ OCR enabled for document processing")
        else:
            self.ocr_processor = None
            logger.info("⚠️  OCR disabled - scanned PDFs and images will fail")
    
    def process_document(self, file_path: Path) -> Optional[Dict[str, Any]]:
        """
        Process a single document and extract geometry content.
        
        Args:
            file_path: Path to document file
        
        Returns:
            Document record dictionary or None if processing fails
        """
        
        if not file_path.exists():
            logger.error(f"❌ File not found: {file_path}")
            return None
        
        file_ext = file_path.suffix.lower()
        if file_ext not in self.SUPPORTED_FORMATS:
            logger.warning(f"⚠️  Unsupported format: {file_ext}")
            return None
        
        logger.info(f"\n{'='*70}")
        logger.info(f"📄 Processing: {file_path.name}")
        logger.info(f"{'='*70}")
        
        try:
            # Extract grade from filename first
            grade_from_filename = self._extract_grade_from_filename(file_path.name)
            
            # Track if OCR was used
            ocr_used = False
            ocr_confidence = None
            ocr_metadata = {}
            
            # ===================================================================
            # EXTRACTION DECISION LOGIC
            # ===================================================================
            
            if file_ext == '.pdf':
                logger.info("📖 File type: PDF")
                text, ocr_metadata = self._extract_pdf(file_path)
                ocr_used = ocr_metadata.get('ocr_used', False)
                ocr_confidence = ocr_metadata.get('avg_confidence')
                
                if ocr_used:
                    logger.info(f"  🔍 OCR was triggered (scanned/image-based PDF)")
                    logger.info(f"  📊 OCR Confidence: {ocr_confidence:.2f}%")
                else:
                    logger.info(f"  📝 Native text extraction used (digital PDF)")
            
            elif file_ext in {'.png', '.jpg', '.jpeg', '.tiff', '.bmp'}:
                logger.info(f"🖼️  File type: Image ({file_ext})")
                logger.info(f"  🔍 OCR will be used (images always require OCR)")
                text, ocr_metadata = self._extract_image(file_path)
                ocr_used = ocr_metadata.get('ocr_used', False)
                ocr_confidence = ocr_metadata.get('confidence')
                
                if ocr_used:
                    logger.info(f"  📊 OCR Confidence: {ocr_confidence:.2f}%")
            
            elif file_ext == '.docx':
                logger.info("📄 File type: DOCX (native extraction)")
                text = self._extract_docx(file_path)
            
            elif file_ext == '.pptx':
                logger.info("📊 File type: PPTX (native extraction)")
                text = self._extract_pptx(file_path)
            
            elif file_ext == '.md':
                logger.info("📝 File type: Markdown (native extraction)")
                text = self._extract_markdown(file_path)
            
            elif file_ext == '.html':
                logger.info("🌐 File type: HTML (native extraction)")
                text = self._extract_html(file_path)
            
            else:  # .txt
                logger.info("📝 File type: Text (native extraction)")
                text = self._extract_text(file_path)
            
            # ===================================================================
            
            # Clean and normalize text
            text = self._clean_text(text)
            
            if not text or not text.strip():
                logger.warning(f"⚠️  No text extracted from {file_path.name}")
                return None
            
            logger.info(f"✅ Extracted {len(text)} characters, {len(text.split())} words")
            
            # Determine source type
            source_type = self._identify_source_type(file_path.name, text)
            logger.info(f"📚 Source type: {source_type}")
            
            # Extract geometry-specific metadata
            metadata = self._extract_geometry_metadata(
                text, 
                file_path, 
                grade_from_filename,
                source_type
            )
            
            logger.info(f"🎓 Grade level: {metadata['grade_level']}")
            logger.info(f"📊 Difficulty: {metadata['difficulty']}")
            logger.info(f"🏷️  Topics: {', '.join(metadata.get('topics', [])) or 'None detected'}")
            
            # Add OCR metadata if OCR was used
            if ocr_used:
                metadata['ocr_processed'] = True
                metadata['ocr_confidence'] = ocr_confidence
                
                # Add additional OCR metadata
                if 'has_diagrams' in ocr_metadata:
                    metadata['ocr_has_diagrams'] = ocr_metadata['has_diagrams']
                if 'has_formulas' in ocr_metadata:
                    metadata['ocr_has_formulas'] = ocr_metadata['has_formulas']
                if 'image_count' in ocr_metadata:
                    metadata['ocr_image_count'] = ocr_metadata['image_count']
            
            # Create document record
            doc_record = {
                'doc_id': self._generate_doc_id(file_path),
                'file_name': file_path.name,
                'file_path': str(file_path),
                'content': text,
                'metadata': metadata,
                'processed_at': datetime.now().isoformat(),
                'char_count': len(text),
                'word_count': len(text.split()),
                'source_type': source_type,
                'ocr_used': ocr_used
            }
            
            # Save processed document
            self._save_processed_document(doc_record)
            
            logger.info(f"✅ Processing complete: {file_path.name}")
            
            return doc_record
        
        except Exception as e:
            logger.error(f"❌ Error processing {file_path}: {e}", exc_info=True)
            return None
    
    def _extract_grade_from_filename(self, filename: str) -> Optional[str]:
        """Extract grade level from filename patterns."""
        filename_lower = filename.lower()
        
        # Pattern 1: classX_Y.pdf or gradeX_Y.pdf
        match = re.search(r'(?:class|grade)(\d+)(?:_|-)?\d*', filename_lower)
        if match:
            grade_num = int(match.group(1))
            if 6 <= grade_num <= 10:
                return f"Grade {grade_num}"
        
        # Pattern 2: Just number like 6.pdf, 7.pdf
        match = re.search(r'^(\d+)(?:_|\.).*\.pdf$', filename_lower)
        if match:
            grade_num = int(match.group(1))
            if 6 <= grade_num <= 10:
                return f"Grade {grade_num}"
        
        return None
    
    def _identify_source_type(self, filename: str, text: str) -> str:
        """Identify the type of source document."""
        filename_lower = filename.lower()
        text_lower = text.lower()
        
        # NCERT identification
        if 'class' in filename_lower and any(f'class{i}' in filename_lower for i in range(6, 11)):
            return 'NCERT Textbook'
        
        # General textbook identification
        if 'textbook' in filename_lower or 'geometry for enjoyment' in text_lower:
            return 'Textbook'
        
        # PPT identification
        if filename.endswith(('.ppt', '.pptx')):
            return 'Presentation'
        
        # Image identification
        if filename.endswith(('.png', '.jpg', '.jpeg', '.tiff', '.bmp')):
            return 'Image/Diagram'
        
        # Default
        return 'Supplementary Material'
    
    def _extract_pdf(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text from PDF with OCR fallback.
        
        ⚠️ OCR TRIGGER POINT:
        - If native text extraction returns empty/whitespace → Calls OCR
        
        Returns:
            Tuple of (text, ocr_metadata)
        """
        text_parts = []
        ocr_metadata = {'ocr_used': False, 'avg_confidence': None}
        
        try:
            logger.info("  📖 Attempting native PDF text extraction...")
            
            # Try native text extraction first using PyPDF
            with open(file_path, 'rb') as file:
                pdf_reader = pypdf.PdfReader(file)
                has_text = False
                
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text and page_text.strip():
                            text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
                            has_text = True
                    except Exception as e:
                        logger.warning(f"    ⚠️  Error extracting page {page_num}: {e}")
            
            # ================================================================
            # OCR DECISION POINT FOR PDFs
            # ================================================================
            if not has_text and self.enable_ocr and self.ocr_processor:
                logger.warning(f"  ⚠️  No native text found - PDF appears to be scanned/image-based")
                logger.info(f"  🔍 Triggering OCR processing...")
                
                ocr_result = self.ocr_processor.process_pdf_with_ocr(
                    str(file_path),
                    extract_images=True,
                    ocr_images=True
                )
                
                text_parts = [ocr_result.get('full_text', '')]
                ocr_metadata = {
                    'ocr_used': True,
                    'avg_confidence': ocr_result.get('avg_confidence', 0.0),
                    'image_count': ocr_result.get('image_count', 0),
                    'has_diagrams': ocr_result.get('has_diagrams', False),
                    'has_formulas': ocr_result.get('has_formulas', False)
                }
                
                if ocr_result.get('avg_confidence', 0) < 70:
                    logger.warning(f"  ⚠️  Low OCR confidence ({ocr_result['avg_confidence']:.1f}%) "
                                 f"for {file_path.name}")
                
            elif not has_text:
                logger.error(f"  ❌ No text found and OCR is disabled!")
            else:
                logger.info(f"  ✅ Native text extraction successful ({len(text_parts)} pages)")
        
        except Exception as e:
            logger.error(f"  ❌ Error extracting PDF {file_path}: {e}")
        
        return '\n\n'.join(text_parts), ocr_metadata
    
    def _extract_image(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """
        Extract text from image file using OCR.
        
        ⚠️ OCR TRIGGER POINT:
        - Image files ALWAYS trigger OCR (no alternative extraction method)
        
        Returns:
            Tuple of (text, ocr_metadata)
        """
        ocr_metadata = {'ocr_used': False, 'confidence': 0.0}
        
        if not self.enable_ocr or not self.ocr_processor:
            logger.error(f"  ❌ OCR disabled - cannot process image {file_path.name}")
            return '', ocr_metadata
        
        try:
            logger.info(f"  🔍 Running OCR on image file...")
            
            result = self.ocr_processor.process_image_file(str(file_path))
            
            text = result.get('text', '')
            confidence = result.get('confidence', 0)
            
            ocr_metadata = {
                'ocr_used': True,
                'confidence': confidence,
                'contains_diagrams': result.get('contains_diagrams', False),
                'contains_formulas': result.get('contains_formulas', False)
            }
            
            if confidence < 70:
                logger.warning(f"  ⚠️  Low OCR confidence ({confidence:.1f}%) for {file_path.name}")
            else:
                logger.info(f"  ✅ OCR successful with {confidence:.1f}% confidence")
            
            return text, ocr_metadata
        
        except Exception as e:
            logger.error(f"  ❌ Error processing image {file_path}: {e}")
            return '', ocr_metadata
    
    def _extract_docx(self, file_path: Path) -> str:
        """Extract text from DOCX file (no OCR needed)."""
        logger.info("  📄 Extracting from DOCX...")
        doc = DocxDocument(file_path)
        
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text_parts.append(row_text)
        
        logger.info(f"    ✅ Extracted {len(text_parts)} paragraphs/tables")
        return '\n\n'.join(text_parts)
    
    def _extract_pptx(self, file_path: Path) -> str:
        """Extract text from PPTX file (no OCR needed)."""
        logger.info("  📊 Extracting from PPTX...")
        prs = Presentation(file_path)
        
        text_parts = []
        ppt_title = file_path.stem.replace('_', ' ').replace('-', ' ').title()
        text_parts.append(f"[Presentation: {ppt_title}]")
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    if shape.text.strip():
                        slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"[Slide {slide_num + 1}]\n" + '\n'.join(slide_text))
        
        logger.info(f"    ✅ Extracted {len(prs.slides)} slides")
        return '\n\n'.join(text_parts)
    
    def _extract_markdown(self, file_path: Path) -> str:
        """Extract text from Markdown file (no OCR needed)."""
        logger.info("  📝 Extracting from Markdown...")
        with open(file_path, 'r', encoding='utf-8') as f:
            md_text = f.read()
        
        html = markdown.markdown(md_text)
        soup = BeautifulSoup(html, 'html.parser')
        
        text = soup.get_text()
        logger.info(f"    ✅ Extracted {len(text)} characters")
        return text
    
    def _extract_html(self, file_path: Path) -> str:
        """Extract text from HTML file (no OCR needed)."""
        logger.info("  🌐 Extracting from HTML...")
        with open(file_path, 'r', encoding='utf-8') as f:
            html_content = f.read()
        
        soup = BeautifulSoup(html_content, 'html.parser')
        
        for element in soup(['script', 'style']):
            element.decompose()
        
        text = soup.get_text()
        logger.info(f"    ✅ Extracted {len(text)} characters")
        return text
    
    def _extract_text(self, file_path: Path) -> str:
        """Extract text from plain text file (no OCR needed)."""
        logger.info("  📝 Extracting from text file...")
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        except UnicodeDecodeError:
            logger.warning("    ⚠️  UTF-8 decode failed, trying latin-1...")
            with open(file_path, 'r', encoding='latin-1') as f:
                text = f.read()
        
        logger.info(f"    ✅ Extracted {len(text)} characters")
        return text
    
    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text."""
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Keep geometry/math symbols
        text = re.sub(r'[^\w\s\-+×÷=°∠∆∏∑√∞≈≠≤≥±∓∝∟⊥∥∦∴∵\(\)\[\]\{\}\/\.\,\:\;]', ' ', text)
        
        # Normalize whitespace
        text = ' '.join(text.split())
        
        return text.strip()
    
    def _extract_geometry_metadata(
        self, 
        text: str, 
        file_path: Path,
        grade_from_filename: Optional[str],
        source_type: str
    ) -> Dict[str, Any]:
        """Extract geometry-specific metadata from text."""
        
        # Determine grade level (prioritize filename, then content)
        if grade_from_filename:
            grade_level = grade_from_filename
        else:
            grade_level = self._infer_grade_level(text, source_type)
        
        metadata = {
            'source_file': file_path.name,
            'source_type': source_type,
            'topics': [],
            'concepts': {},
            'grade_level': grade_level,
            'difficulty': self._assess_difficulty(text, grade_level),
            'chapter_info': self._extract_chapter_info(file_path.name)
        }
        
        # Identify geometry topics
        for topic, pattern in self.GEOMETRY_PATTERNS.items():
            matches = re.findall(pattern, text.lower())
            if matches:
                metadata['concepts'][topic] = list(set(matches))[:10]
                metadata['topics'].append(topic)
        
        # Extract formulas
        formula_pattern = r'[A-Z]\s*=\s*[^.]+(?:\.|$)'
        formulas = re.findall(formula_pattern, text)
        if formulas:
            metadata['formulas'] = formulas[:5]
        
        return metadata
    
    def _extract_chapter_info(self, filename: str) -> Optional[Dict[str, Any]]:
        """Extract chapter information from filename."""
        match = re.search(r'(?:class|grade)(\d+)(?:_|-)(\d+)', filename.lower())
        if match:
            return {
                'grade': int(match.group(1)),
                'chapter': int(match.group(2))
            }
        return None
    
    def _infer_grade_level(self, text: str, source_type: str) -> str:
        """Infer grade level from content complexity."""
        text_lower = text.lower()
        
        if source_type == 'NCERT Textbook':
            return self._infer_ncert_grade(text_lower)
        
        if source_type == 'Textbook':
            return self._infer_textbook_grade(text_lower)
        
        return self._infer_grade_from_keywords(text_lower)
    
    def _infer_ncert_grade(self, text_lower: str) -> str:
        """Infer NCERT grade level from content."""
        grade_scores = {f'Grade {i}': 0 for i in range(6, 11)}
        
        for grade, keywords in self.GRADE_LEVEL_KEYWORDS.items():
            grade_num = grade.split('_')[1]
            grade_label = f'Grade {grade_num}'
            for keyword in keywords:
                if keyword in text_lower:
                    grade_scores[grade_label] += 1
        
        max_grade = max(grade_scores, key=grade_scores.get)
        if grade_scores[max_grade] > 0:
            return max_grade
        
        return 'Grade 6-10 (General)'
    
    def _infer_textbook_grade(self, text_lower: str) -> str:
        """Infer grade for comprehensive textbooks."""
        if any(term in text_lower for term in [
            'points and lines', 'planes', 'segments', 'rays',
            'basic constructions', 'measuring angles'
        ]):
            return 'Grade 6-7'
        
        elif any(term in text_lower for term in [
            'congruent triangles', 'quadrilaterals', 'parallel lines',
            'area formulas', 'volume', 'surface area'
        ]):
            return 'Grade 8-9'
        
        elif any(term in text_lower for term in [
            'coordinate geometry', 'trigonometry', 'similar triangles',
            'circles and tangents', 'proofs of theorems'
        ]):
            return 'Grade 10'
        
        elif any(term in text_lower for term in [
            'analytic geometry', 'conic sections', 'vectors',
            'transformations', 'advanced trigonometry'
        ]):
            return 'Grade 10 (Advanced)'
        
        return 'Grade 6-10 (General)'
    
    def _infer_grade_from_keywords(self, text_lower: str) -> str:
        """Infer grade from keyword analysis."""
        grade_scores = {}
        
        for grade_key, keywords in self.GRADE_LEVEL_KEYWORDS.items():
            grade_num = grade_key.split('_')[1]
            score = sum(1 for keyword in keywords if keyword in text_lower)
            if score > 0:
                grade_scores[f'Grade {grade_num}'] = score
        
        if grade_scores:
            return max(grade_scores, key=grade_scores.get)
        
        # Fallback based on complexity
        complex_terms = ['proof', 'theorem', 'lemma', 'corollary', 'axiom', 'postulate']
        complexity_score = sum(1 for term in complex_terms if term in text_lower)
        
        if complexity_score >= 4:
            return 'Grade 10'
        elif complexity_score >= 2:
            return 'Grade 8-9'
        else:
            return 'Grade 6-7'
    
    def _assess_difficulty(self, text: str, grade_level: str) -> str:
        """Assess content difficulty level based on grade and content."""
        text_lower = text.lower()
        
        complex_terms = ['proof', 'theorem', 'lemma', 'corollary', 'axiom', 
                        'postulate', 'derive', 'prove']
        complexity_score = sum(1 for term in complex_terms if term in text_lower)
        
        if 'Grade 6' in grade_level or 'Grade 7' in grade_level:
            if complexity_score >= 2:
                return 'Intermediate'
            else:
                return 'Beginner'
        
        elif 'Grade 8' in grade_level or 'Grade 9' in grade_level:
            if complexity_score >= 4:
                return 'Advanced'
            elif complexity_score >= 2:
                return 'Intermediate'
            else:
                return 'Beginner'
        
        elif 'Grade 10' in grade_level:
            if complexity_score >= 5:
                return 'Advanced'
            elif complexity_score >= 2:
                return 'Intermediate'
            else:
                return 'Beginner'
        
        else:
            if complexity_score >= 4:
                return 'Advanced'
            elif complexity_score >= 2:
                return 'Intermediate'
            else:
                return 'Beginner'
    
    def _generate_doc_id(self, file_path: Path) -> str:
        """Generate unique document ID."""
        content = f"{file_path.name}{file_path.stat().st_size}{file_path.stat().st_mtime}"
        return hashlib.md5(content.encode()).hexdigest()
    
    def _save_processed_document(self, doc_record: Dict[str, Any]):
        """Save processed document to file."""
        output_path = settings.processed_data_dir / f"{doc_record['doc_id']}.json"
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(doc_record, f, indent=2, ensure_ascii=False)
        
        ocr_info = f" [OCR: ✓]" if doc_record.get('ocr_used') else ""
        logger.info(f"💾 Saved: {doc_record['file_name']} "
                   f"[{doc_record['metadata']['grade_level']}]{ocr_info}")
    
    def process_directory(self, directory: Path) -> List[Dict[str, Any]]:
        """Process all documents in a directory with detailed progress reporting."""
        processed_docs = []
        
        # Get all files
        all_files = [f for f in directory.rglob('*') 
                     if f.is_file() and f.suffix.lower() in self.SUPPORTED_FORMATS]
        
        logger.info(f"\n{'='*70}")
        logger.info(f"📁 Found {len(all_files)} documents to process in: {directory}")
        logger.info(f"{'='*70}\n")
        
        # Track statistics
        ocr_count = 0
        failed_count = 0
        format_counts = {}
        
        for i, file_path in enumerate(all_files, 1):
            logger.info(f"[{i}/{len(all_files)}] Processing: {file_path.name}")
            
            # Track format
            file_ext = file_path.suffix.lower()
            format_counts[file_ext] = format_counts.get(file_ext, 0) + 1
            
            doc_record = self.process_document(file_path)
            
            if doc_record:
                processed_docs.append(doc_record)
                
                if doc_record.get('ocr_used'):
                    ocr_count += 1
            else:
                failed_count += 1
                logger.warning(f"  ✗ Failed to process: {file_path.name}")
        
        # Print comprehensive summary
        self._print_processing_summary(
            all_files, processed_docs, failed_count, ocr_count, format_counts
        )
        
        return processed_docs
    
    def _print_processing_summary(
        self, 
        all_files: List[Path],
        processed_docs: List[Dict[str, Any]],
        failed_count: int,
        ocr_count: int,
        format_counts: Dict[str, int]
    ):
        """Print detailed processing summary."""
        logger.info(f"\n{'='*70}")
        logger.info(f"📊 PROCESSING SUMMARY")
        logger.info(f"{'='*70}")
        logger.info(f"Total files found: {len(all_files)}")
        logger.info(f"Successfully processed: {len(processed_docs)}")
        logger.info(f"Failed: {failed_count}")
        logger.info(f"OCR processed: {ocr_count}")
        
        # Format distribution
        logger.info(f"\n📄 Format Distribution:")
        for fmt, count in sorted(format_counts.items()):
            logger.info(f"  {fmt}: {count} files")
        
        # Grade distribution
        grade_dist = {}
        for doc in processed_docs:
            grade = doc['metadata']['grade_level']
            grade_dist[grade] = grade_dist.get(grade, 0) + 1
        
        logger.info(f"\n🎓 Grade Distribution:")
        for grade, count in sorted(grade_dist.items()):
            logger.info(f"  {grade}: {count} documents")
        
        # Difficulty distribution
        diff_dist = {}
        for doc in processed_docs:
            diff = doc['metadata']['difficulty']
            diff_dist[diff] = diff_dist.get(diff, 0) + 1
        
        logger.info(f"\n📊 Difficulty Distribution:")
        for diff, count in sorted(diff_dist.items()):
            logger.info(f"  {diff}: {count} documents")
        
        # Source type distribution
        source_dist = {}
        for doc in processed_docs:
            source = doc['source_type']
            source_dist[source] = source_dist.get(source, 0) + 1
        
        logger.info(f"\n📚 Source Type Distribution:")
        for source, count in sorted(source_dist.items()):
            logger.info(f"  {source}: {count} documents")
        
        # OCR statistics
        if ocr_count > 0:
            ocr_confidences = [
                doc['metadata'].get('ocr_confidence', 0)
                for doc in processed_docs
                if doc.get('ocr_used') and doc['metadata'].get('ocr_confidence')
            ]
            
            if ocr_confidences:
                avg_conf = sum(ocr_confidences) / len(ocr_confidences)
                logger.info(f"\n🔍 OCR Statistics:")
                logger.info(f"  Documents processed with OCR: {ocr_count}")
                logger.info(f"  Average OCR confidence: {avg_conf:.2f}%")
                logger.info(f"  Min confidence: {min(ocr_confidences):.2f}%")
                logger.info(f"  Max confidence: {max(ocr_confidences):.2f}%")
        
        logger.info(f"{'='*70}\n")
    
    def get_processing_stats(self, processed_docs: List[Dict[str, Any]]) -> Dict[str, Any]:
        """Get detailed processing statistics."""
        stats = {
            'total_documents': len(processed_docs),
            'total_chars': sum(doc['char_count'] for doc in processed_docs),
            'total_words': sum(doc['word_count'] for doc in processed_docs),
            'ocr_processed': sum(1 for doc in processed_docs if doc.get('ocr_used')),
            'grade_distribution': {},
            'difficulty_distribution': {},
            'source_type_distribution': {},
            'topics_distribution': {},
            'avg_ocr_confidence': None
        }
        
        # Calculate distributions
        for doc in processed_docs:
            # Grade
            grade = doc['metadata']['grade_level']
            stats['grade_distribution'][grade] = stats['grade_distribution'].get(grade, 0) + 1
            
            # Difficulty
            diff = doc['metadata']['difficulty']
            stats['difficulty_distribution'][diff] = stats['difficulty_distribution'].get(diff, 0) + 1
            
            # Source type
            source = doc['source_type']
            stats['source_type_distribution'][source] = stats['source_type_distribution'].get(source, 0) + 1
            
            # Topics
            for topic in doc['metadata'].get('topics', []):
                stats['topics_distribution'][topic] = stats['topics_distribution'].get(topic, 0) + 1
        
        # Average OCR confidence
        ocr_confidences = [
            doc['metadata'].get('ocr_confidence', 0) 
            for doc in processed_docs 
            if doc.get('ocr_used') and doc['metadata'].get('ocr_confidence')
        ]
        
        if ocr_confidences:
            stats['avg_ocr_confidence'] = sum(ocr_confidences) / len(ocr_confidences)
        
        return stats
    
    def export_processing_report(
        self, 
        processed_docs: List[Dict[str, Any]], 
        output_path: Path
    ):
        """Export detailed processing report to JSON."""
        stats = self.get_processing_stats(processed_docs)
        
        report = {
            'generated_at': datetime.now().isoformat(),
            'statistics': stats,
            'documents': [
                {
                    'doc_id': doc['doc_id'],
                    'file_name': doc['file_name'],
                    'grade': doc['metadata']['grade_level'],
                    'difficulty': doc['metadata']['difficulty'],
                    'source_type': doc['source_type'],
                    'word_count': doc['word_count'],
                    'ocr_used': doc.get('ocr_used', False),
                    'ocr_confidence': doc['metadata'].get('ocr_confidence'),
                    'topics': doc['metadata'].get('topics', [])
                }
                for doc in processed_docs
            ]
        }
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2, ensure_ascii=False)
        
        logger.info(f"📄 Processing report exported to: {output_path}")


# =============================================================================
# EXAMPLE USAGE AND TESTING
# =============================================================================

if __name__ == "__main__":
    import sys
    
    # Configure logging
    logging.basicConfig(
        level=logging.INFO,
        format='%(asctime)s - %(levelname)s - %(message)s',
        handlers=[
            logging.StreamHandler(sys.stdout),
            logging.FileHandler('document_processing.log')
        ]
    )
    
    print("\n" + "="*70)
    print("GEOMETRY DOCUMENT PROCESSOR WITH OCR - TEST SUITE")
    print("="*70)
    
    # Initialize processor with OCR enabled
    processor = GeometryDocumentProcessor(
        enable_ocr=True,
        ocr_config={
            'min_confidence': 60.0,
            'enable_preprocessing': True,
            # 'tesseract_path': '/usr/local/bin/tesseract'  # Uncomment if needed
        }
    )
    
    # =========================================================================
    # Test 1: Process single PDF (will use native extraction if digital)
    # =========================================================================
    print("\n" + "="*70)
    print("TEST 1: Processing Digital PDF")
    print("="*70)
    
    pdf_path = Path('data/raw/class6_9.pdf')
    if pdf_path.exists():
        doc = processor.process_document(pdf_path)
        if doc:
            print(f"\n✅ Success!")
            print(f"Document ID: {doc['doc_id']}")
            print(f"Grade Level: {doc['metadata']['grade_level']}")
            print(f"Difficulty: {doc['metadata']['difficulty']}")
            print(f"OCR Used: {doc.get('ocr_used', False)}")
            if doc.get('ocr_used'):
                print(f"OCR Confidence: {doc['metadata'].get('ocr_confidence', 0):.2f}%")
    else:
        print(f"⚠️  File not found: {pdf_path}")
    
    # =========================================================================
    # Test 2: Process scanned PDF (will trigger OCR)
    # =========================================================================
    print("\n" + "="*70)
    print("TEST 2: Processing Scanned/Image-based PDF")
    print("="*70)
    
    scanned_pdf = Path('data/raw/scanned_textbook.pdf')
    if scanned_pdf.exists():
        doc = processor.process_document(scanned_pdf)
        if doc:
            print(f"\n✅ Success!")
            print(f"OCR Used: {doc.get('ocr_used', False)}")
            print(f"OCR Confidence: {doc['metadata'].get('ocr_confidence', 0):.2f}%")
            print(f"Word Count: {doc['word_count']}")
    else:
        print(f"⚠️  File not found: {scanned_pdf}")
    
    # =========================================================================
    # Test 3: Process image file (will always trigger OCR)
    # =========================================================================
    print("\n" + "="*70)
    print("TEST 3: Processing Image File")
    print("="*70)
    
    image_path = Path('data/raw/geometry_diagram.png')
    if image_path.exists():
        doc = processor.process_document(image_path)
        if doc:
            print(f"\n✅ Success!")
            print(f"OCR Confidence: {doc['metadata'].get('ocr_confidence', 0):.2f}%")
            print(f"Contains Diagrams: {doc['metadata'].get('ocr_has_diagrams', False)}")
            print(f"Text Preview: {doc['content'][:200]}...")
    else:
        print(f"⚠️  File not found: {image_path}")
    
    # =========================================================================
    # Test 4: Batch process directory
    # =========================================================================
    print("\n" + "="*70)
    print("TEST 4: Batch Processing Directory")
    print("="*70)
    
    data_dir = Path('data/raw')
    if data_dir.exists():
        processed_docs = processor.process_directory(data_dir)
        
        # Export report
        report_path = Path('processing_report.json')
        processor.export_processing_report(processed_docs, report_path)
        
        print(f"\n✅ Batch processing complete!")
        print(f"📄 Report saved to: {report_path}")
    else:
        print(f"⚠️  Directory not found: {data_dir}")
        print("Please create data/raw/ and add geometry documents")
    
    print("\n" + "="*70)
    print("TEST SUITE COMPLETE")
    print("="*70 + "\n")