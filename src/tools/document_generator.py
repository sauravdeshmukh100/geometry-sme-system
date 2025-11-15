#!/usr/bin/env python
"""
Document Generation Tools
Supports PDF, DOCX, and PPT generation for geometry educational content
"""

import os
import logging
from typing import Dict, Any, Optional, List
from datetime import datetime
from io import BytesIO

# PDF Generation
try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, PageBreak,
        Table, TableStyle, Image as RLImage
    )
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

# DOCX Generation
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    PYTHON_DOCX_AVAILABLE = True
except ImportError:
    PYTHON_DOCX_AVAILABLE = False

# PPT Generation
try:
    from pptx import Presentation
    from pptx.util import Inches as PptInches, Pt as PptPt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor as PptRGBColor
    PYTHON_PPTX_AVAILABLE = True
except ImportError:
    PYTHON_PPTX_AVAILABLE = False

logger = logging.getLogger(__name__)


class DocumentGenerationError(Exception):
    """Custom exception for document generation errors."""
    pass


class PDFGenerator:
    """Generate PDF documents for geometry content."""
    
    def __init__(self, output_dir: str = "generated_docs"):
        """
        Initialize PDF generator.
        
        Args:
            output_dir: Directory to save generated PDFs
        """
        if not REPORTLAB_AVAILABLE:
            raise ImportError(
                "ReportLab not installed. Install with: pip install reportlab"
            )
        
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
        
        # Setup styles
        self.styles = getSampleStyleSheet()
        self._setup_custom_styles()
        
        logger.info("PDF Generator initialized")
    
    def _setup_custom_styles(self):
        """Setup custom paragraph styles."""
        # Title style
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#1a73e8'),
            spaceAfter=30,
            alignment=TA_CENTER
        ))
        
        # Heading style
        self.styles.add(ParagraphStyle(
            name='CustomHeading',
            parent=self.styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#202124'),
            spaceAfter=12,
            spaceBefore=12
        ))
        
        # Body style
        self.styles.add(ParagraphStyle(
            name='CustomBody',
            parent=self.styles['BodyText'],
            fontSize=11,
            alignment=TA_JUSTIFY,
            spaceAfter=12
        ))
    
    def generate_quiz_pdf(
        self,
        quiz_data: Dict[str, Any],
        filename: Optional[str] = None,
        include_answers: bool = True
    ) -> str:
        """
        Generate a PDF quiz document.
        
        Args:
            quiz_data: Quiz data from LLM pipeline
            filename: Custom filename (auto-generated if None)
            include_answers: Whether to include answer key
        
        Returns:
            Path to generated PDF file
        """
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            topic = quiz_data.get('topic', 'quiz').replace(' ', '_')
            filename = f"quiz_{topic}_{timestamp}.pdf"
        
        filepath = os.path.join(self.output_dir, filename)
        
        try:
            # Create document
            doc = SimpleDocTemplate(filepath, pagesize=letter)
            story = []
            
            # Title
            title = quiz_data.get('quiz_title', f"Geometry Quiz: {quiz_data.get('topic', 'Unknown')}")
            story.append(Paragraph(title, self.styles['CustomTitle']))
            story.append(Spacer(1, 0.3 * inch))
            
            # Metadata
            grade = quiz_data.get('grade_level', 'N/A')
            date = quiz_data.get('metadata', {}).get('generated_at', datetime.now().strftime('%Y-%m-%d'))
            
            metadata_text = f"<b>Grade Level:</b> {grade} | <b>Date:</b> {date}"
            story.append(Paragraph(metadata_text, self.styles['Normal']))
            story.append(Spacer(1, 0.2 * inch))
            
            # Instructions
            instructions = "Instructions: Answer all questions. Show your work where applicable."
            story.append(Paragraph(instructions, self.styles['CustomBody']))
            story.append(Spacer(1, 0.3 * inch))
            
            # Questions
            questions = quiz_data.get('questions', [])
            
            if not questions:
                # Fallback for freeform quiz
                quiz_text = quiz_data.get('quiz', '')
                story.append(Paragraph(quiz_text.replace('\n', '<br/>'), self.styles['CustomBody']))
            else:
                # Structured quiz
                for q in questions:
                    # Question number and text
                    q_num = q.get('question_number', '?')
                    q_text = q.get('question_text', '')
                    
                    story.append(Paragraph(
                        f"<b>Question {q_num}:</b> {q_text}",
                        self.styles['CustomBody']
                    ))
                    story.append(Spacer(1, 0.1 * inch))
                    
                    # Options (if MCQ)
                    options = q.get('options', [])
                    if options:
                        for opt in options:
                            story.append(Paragraph(f"  {opt}", self.styles['Normal']))
                        story.append(Spacer(1, 0.15 * inch))
                    else:
                        # Short answer space
                        story.append(Paragraph("Answer: _" + "_" * 60, self.styles['Normal']))
                        story.append(Spacer(1, 0.2 * inch))
            
            # Answer key (on separate page)
            if include_answers and questions:
                story.append(PageBreak())
                story.append(Paragraph("Answer Key", self.styles['CustomTitle']))
                story.append(Spacer(1, 0.3 * inch))
                
                for q in questions:
                    q_num = q.get('question_number', '?')
                    answer = q.get('correct_answer', 'N/A')
                    explanation = q.get('explanation', '')
                    
                    story.append(Paragraph(
                        f"<b>Question {q_num}:</b> {answer}",
                        self.styles['CustomBody']
                    ))
                    
                    if explanation:
                        story.append(Paragraph(
                            f"<i>Explanation:</i> {explanation}",
                            self.styles['Normal']
                        ))
                    
                    story.append(Spacer(1, 0.15 * inch))
            
            # Build PDF
            doc.build(story)
            
            logger.info(f"✓ PDF quiz generated: {filepath}")
            return filepath
            
        except Exception as e:
            logger.error(f"Error generating quiz PDF: {e}", exc_info=True)
            raise DocumentGenerationError(f"PDF generation failed: {e}")
    
    def generate_report_pdf(
        self,
        title: str,
        content: str,
        metadata: Optional[Dict[str, Any]] = None,
        filename: Optional[str] = None
    ) -> str:
        """
        Generate a PDF report/document.
        
        Args:
            title: Report title
            content: Main content (supports basic HTML tags)
            metadata: Optional metadata (author, date, etc.)
            filename: Custom filename
        
        Returns:
            Path to generated PDF file
        """
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_title = title.replace(' ', '_')[:30]
            filename = f"report_{safe_title}_{timestamp}.pdf"
        
        filepath = os.path.join(self.output_dir, filename)
        
        try:
            doc = SimpleDocTemplate(filepath, pagesize=letter)
            story = []
            
            # Title
            story.append(Paragraph(title, self.styles['CustomTitle']))
            story.append(Spacer(1, 0.3 * inch))
            
            # Metadata
            if metadata:
                meta_lines = []
                for key, value in metadata.items():
                    meta_lines.append(f"<b>{key.title()}:</b> {value}")
                
                story.append(Paragraph(" | ".join(meta_lines), self.styles['Normal']))
                story.append(Spacer(1, 0.2 * inch))
            
            # Content
            # Split by paragraphs and add them
            paragraphs = content.split('\n\n')
            for para in paragraphs:
                if para.strip():
                    # Replace newlines with breaks
                    para_formatted = para.replace('\n', '<br/>')
                    story.append(Paragraph(para_formatted, self.styles['CustomBody']))
                    story.append(Spacer(1, 0.15 * inch))
            
            # Build PDF
            doc.build(story)
            
            logger.info(f"✓ PDF report generated: {filepath}")
            return filepath
            
        except Exception as e:
            logger.error(f"Error generating report PDF: {e}", exc_info=True)
            raise DocumentGenerationError(f"PDF generation failed: {e}")


class DOCXGenerator:
    """Generate DOCX documents for geometry content."""
    
    def __init__(self, output_dir: str = "generated_docs"):
        """
        Initialize DOCX generator.
        
        Args:
            output_dir: Directory to save generated DOCX files
        """
        if not PYTHON_DOCX_AVAILABLE:
            raise ImportError(
                "python-docx not installed. Install with: pip install python-docx"
            )
        
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
        
        logger.info("DOCX Generator initialized")
    
    def generate_quiz_docx(
        self,
        quiz_data: Dict[str, Any],
        filename: Optional[str] = None,
        include_answers: bool = True
    ) -> str:
        """
        Generate a DOCX quiz document.
        
        Args:
            quiz_data: Quiz data from LLM pipeline
            filename: Custom filename
            include_answers: Whether to include answer key
        
        Returns:
            Path to generated DOCX file
        """
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            topic = quiz_data.get('topic', 'quiz').replace(' ', '_')
            filename = f"quiz_{topic}_{timestamp}.docx"
        
        filepath = os.path.join(self.output_dir, filename)
        
        try:
            doc = Document()
            
            # Title
            title = quiz_data.get('quiz_title', f"Geometry Quiz: {quiz_data.get('topic', 'Unknown')}")
            heading = doc.add_heading(title, level=1)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Metadata
            grade = quiz_data.get('grade_level', 'N/A')
            date = quiz_data.get('metadata', {}).get('generated_at', datetime.now().strftime('%Y-%m-%d'))
            
            meta_para = doc.add_paragraph()
            meta_para.add_run(f"Grade Level: ").bold = True
            meta_para.add_run(f"{grade}  |  ")
            meta_para.add_run("Date: ").bold = True
            meta_para.add_run(date)
            meta_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_paragraph()  # Spacer
            
            # Instructions
            inst = doc.add_paragraph()
            inst.add_run("Instructions: ").bold = True
            inst.add_run("Answer all questions. Show your work where applicable.")
            
            doc.add_paragraph()  # Spacer
            
            # Questions
            questions = quiz_data.get('questions', [])
            
            if not questions:
                # Freeform quiz
                quiz_text = quiz_data.get('quiz', '')
                doc.add_paragraph(quiz_text)
            else:
                # Structured quiz
                for q in questions:
                    q_num = q.get('question_number', '?')
                    q_text = q.get('question_text', '')
                    
                    # Question
                    q_para = doc.add_paragraph()
                    q_para.add_run(f"Question {q_num}: ").bold = True
                    q_para.add_run(q_text)
                    
                    # Options
                    options = q.get('options', [])
                    if options:
                        for opt in options:
                            doc.add_paragraph(opt, style='List Bullet')
                    else:
                        # Answer space
                        doc.add_paragraph("Answer: " + "_" * 60)
                    
                    doc.add_paragraph()  # Spacer
            
            # Answer key
            if include_answers and questions:
                doc.add_page_break()
                
                ans_heading = doc.add_heading("Answer Key", level=1)
                ans_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
                
                doc.add_paragraph()
                
                for q in questions:
                    q_num = q.get('question_number', '?')
                    answer = q.get('correct_answer', 'N/A')
                    explanation = q.get('explanation', '')
                    
                    # Answer
                    ans_para = doc.add_paragraph()
                    ans_para.add_run(f"Question {q_num}: ").bold = True
                    ans_para.add_run(answer)
                    
                    # Explanation
                    if explanation:
                        exp_para = doc.add_paragraph()
                        exp_para.add_run("Explanation: ").italic = True
                        exp_para.add_run(explanation)
                    
                    doc.add_paragraph()
            
            # Save
            doc.save(filepath)
            
            logger.info(f"✓ DOCX quiz generated: {filepath}")
            return filepath
            
        except Exception as e:
            logger.error(f"Error generating quiz DOCX: {e}", exc_info=True)
            raise DocumentGenerationError(f"DOCX generation failed: {e}")
    
    def generate_report_docx(
        self,
        title: str,
        content: str,
        metadata: Optional[Dict[str, Any]] = None,
        filename: Optional[str] = None
    ) -> str:
        """
        Generate a DOCX report/document.
        
        Args:
            title: Report title
            content: Main content
            metadata: Optional metadata
            filename: Custom filename
        
        Returns:
            Path to generated DOCX file
        """
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_title = title.replace(' ', '_')[:30]
            filename = f"report_{safe_title}_{timestamp}.docx"
        
        filepath = os.path.join(self.output_dir, filename)
        
        try:
            doc = Document()
            
            # Title
            heading = doc.add_heading(title, level=1)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Metadata
            if metadata:
                meta_para = doc.add_paragraph()
                meta_items = [f"{k.title()}: {v}" for k, v in metadata.items()]
                meta_para.add_run(" | ".join(meta_items))
                meta_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_paragraph()
            
            # Content
            paragraphs = content.split('\n\n')
            for para_text in paragraphs:
                if para_text.strip():
                    doc.add_paragraph(para_text.strip())
            
            # Save
            doc.save(filepath)
            
            logger.info(f"✓ DOCX report generated: {filepath}")
            return filepath
            
        except Exception as e:
            logger.error(f"Error generating report DOCX: {e}", exc_info=True)
            raise DocumentGenerationError(f"DOCX generation failed: {e}")


class PPTGenerator:
    """Generate PowerPoint presentations for geometry content."""
    
    def __init__(self, output_dir: str = "generated_docs"):
        """
        Initialize PPT generator.
        
        Args:
            output_dir: Directory to save generated PPT files
        """
        if not PYTHON_PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx not installed. Install with: pip install python-pptx"
            )
        
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
        
        logger.info("PPT Generator initialized")
    
    def generate_concept_ppt(
        self,
        concept: str,
        explanation: str,
        examples: Optional[List[str]] = None,
        metadata: Optional[Dict[str, Any]] = None,
        filename: Optional[str] = None
    ) -> str:
        """
        Generate a PowerPoint presentation explaining a concept.
        
        Args:
            concept: Concept name
            explanation: Detailed explanation
            examples: Optional list of examples
            metadata: Optional metadata
            filename: Custom filename
        
        Returns:
            Path to generated PPT file
        """
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            safe_concept = concept.replace(' ', '_')[:30]
            filename = f"presentation_{safe_concept}_{timestamp}.pptx"
        
        filepath = os.path.join(self.output_dir, filename)
        
        try:
            prs = Presentation()
            prs.slide_width = PptInches(10)
            prs.slide_height = PptInches(7.5)
            
            # Slide 1: Title
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = concept
            
            if metadata:
                grade = metadata.get('grade_level', '')
                subtitle.text = f"{grade}\n{datetime.now().strftime('%B %d, %Y')}"
            
            # Slide 2: Overview
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            title = slide.shapes.title
            body = slide.placeholders[1]
            
            title.text = "Overview"
            
            # Split explanation into bullet points
            paragraphs = explanation.split('\n\n')
            tf = body.text_frame
            
            for i, para in enumerate(paragraphs[:5]):  # Max 5 points
                if i == 0:
                    tf.text = para.strip()
                else:
                    p = tf.add_paragraph()
                    p.text = para.strip()
                    p.level = 0
            
            # Slide 3+: Examples (if provided)
            if examples:
                for i, example in enumerate(examples[:3], 1):  # Max 3 examples
                    slide = prs.slides.add_slide(bullet_slide_layout)
                    
                    title = slide.shapes.title
                    body = slide.placeholders[1]
                    
                    title.text = f"Example {i}"
                    body.text = example
            
            # Save
            prs.save(filepath)
            
            logger.info(f"✓ PPT presentation generated: {filepath}")
            return filepath
            
        except Exception as e:
            logger.error(f"Error generating PPT: {e}", exc_info=True)
            raise DocumentGenerationError(f"PPT generation failed: {e}")
    
    def generate_quiz_ppt(
        self,
        quiz_data: Dict[str, Any],
        filename: Optional[str] = None
    ) -> str:
        """
        Generate a PowerPoint quiz presentation.
        
        Args:
            quiz_data: Quiz data from LLM pipeline
            filename: Custom filename
        
        Returns:
            Path to generated PPT file
        """
        if filename is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            topic = quiz_data.get('topic', 'quiz').replace(' ', '_')
            filename = f"quiz_{topic}_{timestamp}.pptx"
        
        filepath = os.path.join(self.output_dir, filename)
        
        try:
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            quiz_title = quiz_data.get('quiz_title', f"Quiz: {quiz_data.get('topic', 'Geometry')}")
            title.text = quiz_title
            
            grade = quiz_data.get('grade_level', 'N/A')
            subtitle.text = f"Grade Level: {grade}"
            
            # Question slides
            questions = quiz_data.get('questions', [])
            
            for q in questions:
                # Question slide
                blank_slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Question number and text
                left = PptInches(1)
                top = PptInches(1)
                width = PptInches(8)
                height = PptInches(1)
                
                q_num = q.get('question_number', '?')
                q_text = q.get('question_text', '')
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = f"Question {q_num}"
                
                # Format question number
                p = tf.paragraphs[0]
                p.font.size = PptPt(28)
                p.font.bold = True
                
                # Question text
                left = PptInches(1)
                top = PptInches(2)
                height = PptInches(1.5)
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = q_text
                tf.word_wrap = True
                
                # Options
                options = q.get('options', [])
                if options:
                    left = PptInches(1.5)
                    top = PptInches(3.5)
                    
                    for opt in options:
                        height = PptInches(0.5)
                        txBox = slide.shapes.add_textbox(left, top, width - 1, height)
                        tf = txBox.text_frame
                        tf.text = opt
                        top += PptInches(0.6)
            
            # Save
            prs.save(filepath)
            
            logger.info(f"✓ PPT quiz generated: {filepath}")
            return filepath
            
        except Exception as e:
            logger.error(f"Error generating quiz PPT: {e}", exc_info=True)
            raise DocumentGenerationError(f"PPT generation failed: {e}")


class DocumentGenerator:
    """
    Unified document generator supporting PDF, DOCX, and PPT.
    Provides a single interface for all document generation tasks.
    """
    
    def __init__(self, output_dir: str = "generated_docs"):
        """
        Initialize document generator.
        
        Args:
            output_dir: Directory for generated documents
        """
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)
        
        # Initialize generators
        self.pdf = PDFGenerator(output_dir) if REPORTLAB_AVAILABLE else None
        self.docx = DOCXGenerator(output_dir) if PYTHON_DOCX_AVAILABLE else None
        self.ppt = PPTGenerator(output_dir) if PYTHON_PPTX_AVAILABLE else None
        
        logger.info("Document Generator initialized")
        logger.info(f"  PDF: {'✓' if self.pdf else '✗'}")
        logger.info(f"  DOCX: {'✓' if self.docx else '✗'}")
        logger.info(f"  PPT: {'✓' if self.ppt else '✗'}")
    
    def generate(
        self,
        content_type: str,
        format: str,
        data: Dict[str, Any],
        filename: Optional[str] = None,
        **kwargs
    ) -> str:
        """
        Generate document in specified format.
        
        Args:
            content_type: Type of content ('quiz', 'report', 'concept', 'explanation')
            format: Output format ('pdf', 'docx', 'pptx')
            data: Content data
            filename: Custom filename
            **kwargs: Additional format-specific options
        
        Returns:
            Path to generated document
        
        Raises:
            DocumentGenerationError: If generation fails
        """
        format = format.lower()
        content_type = content_type.lower()
        
        try:
            if format == 'pdf':
                if not self.pdf:
                    raise DocumentGenerationError("PDF generation not available. Install reportlab.")
                
                if content_type == 'quiz':
                    return self.pdf.generate_quiz_pdf(data, filename, **kwargs)
                elif content_type in ['report', 'explanation', 'concept']:
                    title = data.get('concept', data.get('title', 'Report'))
                    content = data.get('explanation', data.get('content', ''))
                    metadata = data.get('metadata', {})
                    return self.pdf.generate_report_pdf(title, content, metadata, filename)
                
            elif format == 'docx':
                if not self.docx:
                    raise DocumentGenerationError("DOCX generation not available. Install python-docx.")
                
                if content_type == 'quiz':
                    return self.docx.generate_quiz_docx(data, filename, **kwargs)
                elif content_type in ['report', 'explanation', 'concept']:
                    title = data.get('concept', data.get('title', 'Report'))
                    content = data.get('explanation', data.get('content', ''))
                    metadata = data.get('metadata', {})
                    return self.docx.generate_report_docx(title, content, metadata, filename)
                
            elif format in ['pptx', 'ppt']:
                if not self.ppt:
                    raise DocumentGenerationError("PPT generation not available. Install python-pptx.")
                
                if content_type == 'quiz':
                    return self.ppt.generate_quiz_ppt(data, filename)
                elif content_type in ['concept', 'explanation']:
                    concept = data.get('concept', 'Concept')
                    explanation = data.get('explanation', '')
                    examples = data.get('examples', [])
                    metadata = data.get('metadata', {})
                    return self.ppt.generate_concept_ppt(concept, explanation, examples, metadata, filename)
            
            else:
                raise DocumentGenerationError(f"Unsupported format: {format}")
            
            raise DocumentGenerationError(f"Unsupported content_type: {content_type}")
            
        except DocumentGenerationError:
            raise
        except Exception as e:
            logger.error(f"Document generation error: {e}", exc_info=True)
            raise DocumentGenerationError(f"Generation failed: {e}")
    
    def get_supported_formats(self) -> Dict[str, bool]:
        """Get dictionary of supported formats."""
        return {
            'pdf': self.pdf is not None,
            'docx': self.docx is not None,
            'pptx': self.ppt is not None
        }